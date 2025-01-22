#!/usr/bin/env python3
# -*- coding: utf-8 -*-
# Auteur : Mickael Dorigny @ IT-Connect.fr

# Importation des modules nécessaires
from pptx import Presentation  # Pour manipuler des fichiers PowerPoint (.pptx)
import argparse  # Pour gérer les arguments passés en ligne de commande

# Fonction pour ajouter une diapositive à une présentation existante
def add_slide(prs: object, layout_id: int, placeholders: dict) -> object:
    """
    Ajoute une diapositive à la présentation.

    Arguments :
    prs -- Objet Presentation représentant la présentation en cours
    layout_id -- Identifiant du modèle de mise en page à appliquer à la nouvelle diapositive

    Retourne :
    prs -- Objet Presentation avec la diapositive ajoutée
    """
    layout = prs.slide_layouts[layout_id]  # Sélection du modèle de mise en page
    slide = prs.slides.add_slide(layout)  # Ajout de la diapositive avec le modèle sélectionné

    if placeholders is not None:
        for key, value in placeholders.items():
            placeholder = slide.placeholders[key]
            placeholder.text = value

    return prs

# Fonction principale exécutant le script
def main(args: object) -> None:
    """
    Point d'entrée principal du script.

    Arguments :
    args -- Arguments passés en ligne de commande
    """
    # Création de la présentation à partir d'un fichier .pptx existant
    presentation = Presentation(args.pptx_file)

    # Ajout de diapositives à la présentation
    placeholders = {
		0 : "Mon Titre",
		1 : "Mon sous-titre"
	}
    presentation = add_slide(presentation, 0, placeholders)  # Première diapositive avec le modèle 0

    placeholders = {
        0 : "Titre slide 1",
        1 : "Mon contenu"
    }
    presentation = add_slide(presentation, 1, placeholders)  # Deuxième diapositive avec le modèle 1
    
    placeholders = {
        0 : "Titre slide 2",
        1 : "Mon autre contenu"
    }
    presentation = add_slide(presentation, 1, placeholders)  # Troisième diapositive avec le modèle 1

    # Sauvegarde de la présentation modifiée dans un nouveau fichier
    presentation.save("Nouvelle_presentation.pptx")

# Point d'entrée du script
if __name__ == '__main__':
    # Configuration des arguments à passer en ligne de commande
    parser = argparse.ArgumentParser(prog='PROG')  # Initialisation du parser
    parser.add_argument("-f", '--pptx-file', help='Path to the PPTX template', required=True)  
    # Argument obligatoire spécifiant le chemin du fichier .pptx

    args = parser.parse_args()  # Analyse des arguments fournis
    main(args)  # Appel de la fonction principale


