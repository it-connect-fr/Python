#!/usr/bin/env python3
# -*- coding: utf-8 -*-
from pptx import Presentation
import argparse


def show_layout_id(presentation:object) -> None:
	"""
	Displays the slide layouts and their corresponding placeholder IDs in a presentation.

	This function iterates through all slide layouts in a presentation (using the 'prs' object),
	printing the layout ID and name, followed by the placeholder IDs and names for each layout. The function
	is useful for debugging and understanding the structure of slide layouts and placeholders in a presentation.

	No arguments and no return value.
	"""
	
	# Initialize layout ID counter
	layout_id = 0
	
	# Iterate over all slide layouts in the presentation
	for layout in presentation.slide_layouts:
		print(f"{layout_id} - {layout.name}")  
		placeholder_id = 0
		
		# Iterate over all placeholders in the current layout
		for placeholder in layout.placeholders:
			print(f"\t{placeholder.placeholder_format.idx} - {placeholder.name}") 
			placeholder_id += 1
		
		layout_id += 1  # Increment layout ID counter
	
def main(args:object) -> None:
	presentation = Presentation(args.pptx_file)
	show_layout_id(presentation)

if __name__ == '__main__':
	# create the top-level parser
	parser = argparse.ArgumentParser(prog='PROG')
	parser.add_argument("-f", '--pptx-file', help='Path to the PPTX template', required=True)
	
	args = parser.parse_args()
	main(args)
