#!/usr/bin/env python3
# -*- coding: utf-8 -*-
# Auteur : Mickael Dorigny @ IT-Connect.fr

# Importation des modules nécessaires
from pptx import Presentation  # Pour manipuler des fichiers PowerPoint (.pptx)
import argparse  # Pour gérer les arguments passés en ligne de commande

# Fonction pour ajouter une diapositive à une présentation existante
def add_slide(prs: object, layout_id: int, placeholders: dict) -> object:
    """
    Ajoute une diapositive à la présentation et remplit les zones réservées (placeholders) si elles sont spécifiées.

    Arguments :
    prs -- Objet Presentation représentant la présentation en cours
    layout_id -- Identifiant du modèle de mise en page à appliquer à la nouvelle diapositive
    placeholders -- Dictionnaire contenant les indices des zones réservées et leurs textes associés

    Retourne :
    prs -- Objet Presentation avec la diapositive ajoutée
    """
    layout = prs.slide_layouts[layout_id]  # Sélection du modèle de mise en page
    slide = prs.slides.add_slide(layout)  # Ajout de la diapositive avec le modèle sélectionné

    # Remplir les placeholders si un dictionnaire a été fourni
    if placeholders is not None:
        for key, value in placeholders.items():
            placeholder = slide.placeholders[key]  # Accès à une zone réservée par son index
            if isinstance(value, str) and value.startswith("/") and value.endswith(".png"):
                placeholder.insert_picture(value)
            else :
                placeholder.text = value  # Assignation du texte au placeholder

    return prs

# Fonction principale exécutant le script
def main(args: object) -> None:
    """
    Point d'entrée principal du script.

    Arguments :
    args -- Arguments passés en ligne de commande
    """
    # Création de la présentation à partir d'un fichier .pptx existant
    presentation = Presentation(args.pptx_file)

    # Ajout de diapositives avec des textes spécifiques dans les placeholders
    placeholders = {
        0: "Mon Titre",  # Texte pour le placeholder 0 (souvent le titre)
        1: "Mon sous-titre"  # Texte pour le placeholder 1 (souvent le sous-titre)
    }
    presentation = add_slide(presentation, 0, placeholders)  # Première diapositive avec le modèle 0

    placeholders = {
        0: "Titre slide 1",  # Titre de la première diapositive ajoutée
        1: "Mon contenu"  # Contenu principal de cette diapositive
    }
    presentation = add_slide(presentation, 1, placeholders)  # Deuxième diapositive avec le modèle 1
    
    placeholders = {
        0: "Mon image",  
        1: "/tmp/image.png",
        2: "Logo IT-Connect"  
    }
    presentation = add_slide(presentation, 8, placeholders)  # Troisième diapositive avec le modèle 1

    # Sauvegarde de la présentation modifiée dans un nouveau fichier
    presentation.save("Nouvelle_presentation.pptx")

# Point d'entrée du script
if __name__ == '__main__':
    # Configuration des arguments à passer en ligne de commande
    parser = argparse.ArgumentParser(prog='PROG')  # Initialisation du parser
    parser.add_argument("-f", '--pptx-file', help='Path to the PPTX template', required=True)  
    # Argument obligatoire spécifiant le chemin du fichier .pptx

    args = parser.parse_args()  # Analyse des arguments fournis
    main(args)  # Appel de la fonction principale
